from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette (dark tech / cyan) ─────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # near-black
BG_CARD  = RGBColor(0x16, 0x1B, 0x27)   # card bg
BG_CARD2 = RGBColor(0x1E, 0x29, 0x3B)   # lighter card
CYAN     = RGBColor(0x06, 0xB6, 0xD4)   # primary accent
TEAL     = RGBColor(0x0D, 0x94, 0x88)   # secondary accent
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x94, 0xA3, 0xB8)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
RED      = RGBColor(0xEF, 0x44, 0x44)
YELLOW   = RGBColor(0xEA, 0xB3, 0x08)

W = Inches(13.33)
H = Inches(7.5)
TOTAL = 13

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ───────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = BG
    return s

def tb(s, text, l, t, w, h, size=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def blist(s, items, l, t, w, h, size=16, color=GRAY, marker="  •  ", gap=5):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(gap)
        r = p.add_run()
        r.text = marker + item
        r.font.size = Pt(size)
        r.font.color.rgb = color

def bar(s, color=CYAN, top=Inches(0.07)):
    b = s.shapes.add_shape(1, 0, top, W, Inches(0.055))
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background()

def tag(s, text, color=CYAN):
    tb(s, text.upper(), Inches(0.45), Inches(0.22), Inches(8), Inches(0.38),
       size=10, bold=True, color=color)

def snum(s, n):
    tb(s, f"{n} / {TOTAL}", Inches(12.1), Inches(7.1), Inches(1.1), Inches(0.35),
       size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def card(s, l, t, w, h, border=BG_CARD2, fill=BG_CARD):
    c = s.shapes.add_textbox(l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border
    return c

def stat_box(s, value, label, l, t, color=CYAN):
    tb(s, value, l, t, Inches(3.0), Inches(0.85),
       size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(s, label, l, t + Inches(0.82), Inches(3.0), Inches(0.45),
       size=13, color=GRAY, align=PP_ALIGN.CENTER)

def level_badge(s, num, label, color, l, t):
    b = s.shapes.add_textbox(l, t, Inches(1.6), Inches(0.5))
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background()
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"LEVEL {num}  —  {label}"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = BG

def add_notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# ─────────────────────────────────────────────────────────────────────────────
# S1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, CYAN)

tb(s, "Claude Code", Inches(0.7), Inches(1.2), W - Inches(1.4), Inches(1.6),
   size=72, bold=True, align=PP_ALIGN.CENTER)
tb(s, "From Zero to Mastery", Inches(0.7), Inches(2.75), W - Inches(1.4), Inches(0.9),
   size=36, color=CYAN, align=PP_ALIGN.CENTER)
tb(s, "6 Levels to Becoming a Claude Code Expert",
   Inches(0.7), Inches(3.65), W - Inches(1.4), Inches(0.6),
   size=20, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# 6 level dots
colors6 = [CYAN, TEAL, GREEN, YELLOW, ORANGE, RED]
labels6  = ["L1\nCMD", "L2\nPLAN", "L3\nCTX", "L4\nMCP", "L5\nSKILL", "L6\nAGENT"]
x = Inches(2.8)
for col, lbl in zip(colors6, labels6):
    c = s.shapes.add_textbox(x, Inches(4.6), Inches(1.3), Inches(0.9))
    c.fill.solid(); c.fill.fore_color.rgb = col
    c.line.fill.background()
    tf = c.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lbl
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BG
    x += Inches(1.35)

tb(s, "Tech Talk  |  2026", Inches(0.7), Inches(6.6), W - Inches(1.4), Inches(0.5),
   size=13, color=GRAY, align=PP_ALIGN.CENTER)
snum(s, 1)
add_notes(s, """Welcome everyone. Today I'm going to show you something that fundamentally changes how you build with AI.
Claude Code is not a chatbot. It's a terminal-based AI coding partner — and by the end of this talk you'll understand exactly how to go from a complete beginner to running multi-agent teams.
There are 6 levels. Most people never get past Level 1. Let's change that.
[Pause — let the audience look at the 6 level badges]""")

# ─────────────────────────────────────────────────────────────────────────────
# S2 — What is Claude Code?
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, CYAN)
tag(s, "Introduction")

tb(s, "What is Claude Code?", Inches(0.6), Inches(0.85), Inches(12), Inches(0.75),
   size=38, bold=True)
tb(s, "It's not a chatbot. It lives inside your project.",
   Inches(0.6), Inches(1.65), Inches(12), Inches(0.55), size=21, color=CYAN)

# Left: what it does
blist(s, [
    "Anthropic's terminal-based AI coding tool",
    "Reads your files — understands full project context",
    "Writes code, runs it, fixes errors, deploys",
    "You don't need to be a developer",
    "You just need to know WHAT you want",
], Inches(0.6), Inches(2.5), Inches(6.0), Inches(3.5), size=18)

# Right: comparison cards
tb(s, "Regular AI Chat", Inches(7.0), Inches(2.4), Inches(5.5), Inches(0.45),
   size=13, bold=True, color=RED)
c1 = card(s, Inches(7.0), Inches(2.88), Inches(5.5), Inches(0.85), border=RED)
tb_c = c1.text_frame.paragraphs[0]
r = tb_c.add_run(); r.text = "  Input  →  Output  (you copy-paste manually)"
r.font.size = Pt(13); r.font.color.rgb = GRAY

tb(s, "Claude Code", Inches(7.0), Inches(4.0), Inches(5.5), Inches(0.45),
   size=13, bold=True, color=GREEN)
c2 = card(s, Inches(7.0), Inches(4.48), Inches(5.5), Inches(1.55), border=GREEN)
tf2 = c2.text_frame; tf2.word_wrap = True
steps = ["Reads project  →  Writes  →  Runs", "Fixes errors  →  Tests  →  Deploys"]
for i, step in enumerate(steps):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run(); r.text = "  " + step
    r.font.size = Pt(13); r.font.color.rgb = WHITE

snum(s, 2)
add_notes(s, """The most important thing I can tell you about Claude Code is what it is NOT.
It is not Claude.ai in a browser. It is not copy-paste AI.
Claude Code lives INSIDE your terminal. It reads your actual files. It runs your actual commands. It sees your git history.
The difference in output quality between 'chat AI' and Claude Code is like the difference between asking a stranger for directions and having a navigator in your car.
You don't need to be a programmer to use it. You need to know what outcome you want and describe it clearly.""")

# ─────────────────────────────────────────────────────────────────────────────
# S3 — How is Claude Code Different?
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, CYAN)
tag(s, "Comparison")

tb(s, "Claude Code vs. Everything Else",
   Inches(0.6), Inches(0.85), Inches(12), Inches(0.75), size=36, bold=True)

# Comparison table
headers = ["Feature", "ChatGPT / Gemini", "Claude Code"]
col_w   = [Inches(3.8), Inches(3.8), Inches(3.8)]
col_x   = [Inches(0.5), Inches(4.4), Inches(8.3)]
row_h   = Inches(0.7)
rows = [
    ("Reads your files",   "No",  "Yes"),
    ("Runs commands",      "No",  "Yes"),
    ("Manages git",        "No",  "Yes"),
    ("Fixes own errors",   "No",  "Yes"),
    ("Full project context","No", "Yes"),
    ("Deploys code",       "No",  "Yes"),
]
# Header row
for i, (hdr, cx) in enumerate(zip(headers, col_x)):
    c = s.shapes.add_textbox(cx, Inches(1.75), col_w[i], Inches(0.5))
    c.fill.solid()
    c.fill.fore_color.rgb = CYAN if i == 2 else BG_CARD2
    c.line.fill.background()
    tf = c.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = hdr
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = BG if i == 2 else WHITE

# Data rows
y = Inches(2.3)
for row in rows:
    for i, (val, cx) in enumerate(zip(row, col_x)):
        c = s.shapes.add_textbox(cx, y, col_w[i], row_h - Inches(0.05))
        c.fill.solid()
        c.fill.fore_color.rgb = BG_CARD
        c.line.color.rgb = BG_CARD2
        tf = c.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(14)
        r.font.color.rgb = (GREEN if val == "Yes" else RED if val == "No" else WHITE)
        r.font.bold = (val in ("Yes", "No"))
    y += row_h

# Big stat
stat_box(s, "90%", "of users stuck at Level 1", Inches(9.1), Inches(5.2), RED)

snum(s, 3)
add_notes(s, """Let's be concrete about what makes Claude Code different.
Every row in this table is a capability that ChatGPT, Gemini, and regular Claude simply don't have.
'Reads your files' — this sounds small but it's everything. Claude Code knows your codebase. You don't explain it every session.
'Fixes own errors' — when a script crashes, Claude Code reads the traceback, understands the cause, and fixes it. You don't babysit it.
And here's the uncomfortable stat: 90% of users are stuck at Level 1. They treat Claude Code like a command executor. They get average results.
The next slides will show you how to move up the levels.""")

# ─────────────────────────────────────────────────────────────────────────────
# S4 — Setup & Cost
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, TEAL)
tag(s, "Setup")

tb(s, "Setup & Cost Overview",
   Inches(0.6), Inches(0.85), Inches(12), Inches(0.7), size=36, bold=True)

# 3 setup methods
methods = [
    ("Terminal", "npm install -g\n@anthropic-ai/claude-code", CYAN),
    ("VS Code",  "Install the\nClaude Code extension", TEAL),
    ("Desktop",  "Download the\nClaude desktop app", GREEN),
]
x = Inches(0.5)
for title, desc, col in methods:
    c = card(s, x, Inches(1.75), Inches(3.8), Inches(1.4), border=col)
    tf = c.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = f"  {title}"
    r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = col
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    r2 = p2.add_run(); r2.text = f"  {desc}"
    r2.font.size = Pt(13); r2.font.color.rgb = GRAY
    x += Inches(4.3)

tb(s, "Recommended: Terminal or VS Code",
   Inches(0.5), Inches(3.3), Inches(12), Inches(0.4),
   size=13, color=YELLOW, italic=True)

# Plans
plans = [
    ("Pro",  "$20 / mo",  "Basic access",               GRAY),
    ("Max",  "$100 / mo", "5x-20x tokens + Opus 4.6",   CYAN),
    ("API",  "Pay-as-go", "Full control, production use", ORANGE),
]
x = Inches(0.5)
for name, price, desc, col in plans:
    c = card(s, x, Inches(3.8), Inches(3.8), Inches(1.3), border=col)
    tf = c.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = f"  {name} — {price}"
    r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = col
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    r2 = p2.add_run(); r2.text = f"  {desc}"
    r2.font.size = Pt(13); r2.font.color.rgb = GRAY
    x += Inches(4.3)

# Key commands
tb(s, "Key commands:", Inches(0.5), Inches(5.3), Inches(5), Inches(0.4),
   size=13, bold=True, color=WHITE)
cmds = "/plan  /clear  /compact  /context  /cost  /model  /init"
tb(s, cmds, Inches(0.5), Inches(5.72), Inches(12), Inches(0.5),
   size=14, color=CYAN, bold=True)

# Models
tb(s, "Sonnet 4.6  — 80% of daily work        |        Opus 4.6  — complex architecture & hard bugs",
   Inches(0.5), Inches(6.35), Inches(12), Inches(0.5),
   size=14, color=GRAY, align=PP_ALIGN.CENTER)

snum(s, 4)
add_notes(s, """Let's get practical. Three ways to install Claude Code — Terminal, VS Code extension, or Desktop app.
I recommend Terminal or VS Code. You want to be close to your files.
For pricing: Pro at $20 is fine to start. Max at $100 gives you 5 to 20 times more tokens and access to Opus 4.6 — worth it when you're doing serious work.
API is pay-as-you-go — best for production pipelines.
Two models matter: Sonnet 4.6 handles 80% of everything. Switch to Opus 4.6 only when you're tackling complex architecture or really hard bugs — it's significantly more capable but costs more.
Learn these commands by heart: /plan before big tasks, /clear to reset context, /compact when the session is getting heavy.""")

# ─────────────────────────────────────────────────────────────────────────────
# S5 — Level 1: The Commander
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, CYAN)
tag(s, "Level 1")
level_badge(s, 1, "THE COMMANDER", CYAN, Inches(0.5), Inches(0.18))

tb(s, "Level 1: The Commander",
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75), size=36, bold=True)

stat_box(s, "90%", "of users live here", Inches(10.0), Inches(0.7), RED)

# What it looks like
c = card(s, Inches(0.5), Inches(1.8), Inches(5.9), Inches(2.2), border=RED)
tf = c.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = '  "Build me a landing page"'
r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.space_before = Pt(8)
r2 = p2.add_run(); r2.text = "  Result: generic purple/blue AI-looking site"
r2.font.size = Pt(15); r2.font.color.rgb = RED

blist(s, [
    "No context given  →  Claude guesses  →  low quality",
    "Output targets 'the average' of all similar requests",
    "User gets frustrated, blames the AI",
    "Real problem: the instruction, not the model",
], Inches(0.5), Inches(4.2), Inches(5.9), Inches(2.8), size=16)

# Rule callout
c2 = card(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(3.5), border=YELLOW, fill=BG_CARD2)
tb(s, "The Rule", Inches(7.0), Inches(1.95), Inches(5.6), Inches(0.5),
   size=14, bold=True, color=YELLOW)
tb(s, '"No context =\nClaude targets\nthe average."',
   Inches(7.0), Inches(2.5), Inches(5.6), Inches(2.0),
   size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

snum(s, 5)
add_notes(s, """Level 1 is where 90% of people live — and it's where most frustration with AI comes from.
You type 'build me a landing page'. Claude builds something. It looks generic. You think the AI is bad.
The AI is not bad. YOUR INSTRUCTION is bad.
When you give no context, Claude has to guess. And when Claude guesses, it targets the statistical average of every similar request it has ever seen.
Average prompt = average output. This is not an AI problem. This is a communication problem.
The key insight of Level 1 is simple: no context = Claude targets the average. You must give context. That's it. That's the entire unlock.""")

# ─────────────────────────────────────────────────────────────────────────────
# S6 — Level 2: The Planner
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, TEAL)
tag(s, "Level 2")
level_badge(s, 2, "THE PLANNER", TEAL, Inches(0.5), Inches(0.18))

tb(s, "Level 2: The Planner",
   Inches(0.6), Inches(0.85), Inches(12), Inches(0.75), size=36, bold=True)
tb(s, "Use /plan before you build anything.",
   Inches(0.6), Inches(1.62), Inches(12), Inches(0.55), size=20, color=TEAL)

blist(s, [
    "Type /plan — Claude becomes your project partner",
    "It asks: language? structure? theme? technology?",
    "Even basic /plan use dramatically improves output quality",
    "Ask Claude to challenge your ideas — devil's advocate mode",
    "Claude is a collaborator. Not a command executor.",
], Inches(0.5), Inches(2.4), Inches(6.5), Inches(3.5), size=17)

# Plan flow
flow = [
    ("/plan", TEAL),
    ("Claude\nasks\nquestions", CYAN),
    ("You\nanswer", WHITE),
    ("Aligned\noutput", GREEN),
]
x = Inches(7.2)
for i, (label, col) in enumerate(flow):
    c = s.shapes.add_textbox(x, Inches(2.8), Inches(1.3), Inches(1.3))
    c.fill.solid(); c.fill.fore_color.rgb = BG_CARD
    c.line.color.rgb = col
    tf = c.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = col
    if i < 3:
        tb(s, "→", x + Inches(1.3), Inches(3.2), Inches(0.5), Inches(0.5),
           size=20, color=GRAY, align=PP_ALIGN.CENTER)
    x += Inches(1.55)

# Devil's advocate tip
c2 = card(s, Inches(7.0), Inches(4.5), Inches(5.9), Inches(1.7), border=YELLOW)
tb(s, "Pro tip: Ask Claude to challenge you",
   Inches(7.2), Inches(4.65), Inches(5.5), Inches(0.45),
   size=13, bold=True, color=YELLOW)
tb(s, '"What are the weaknesses in this approach?\nWhat would you do differently?"',
   Inches(7.2), Inches(5.15), Inches(5.5), Inches(0.9),
   size=14, color=WHITE, italic=True)

snum(s, 6)
add_notes(s, """Level 2 is about switching from commander to collaborator.
Before you build anything — type /plan. That's it. That's the entire unlock for Level 2.
When you type /plan, Claude stops executing and starts asking. What language are you using? What's the target audience? What does success look like?
This one habit — /plan before building — will improve your output quality more than any prompt engineering trick.
The mindset shift: Claude is not a tool waiting for commands. Claude is a project partner. Treat it like one.
And here's a power move: ask Claude to challenge your ideas. 'What are the weaknesses in this approach? What would you do differently?' You'll get better solutions than you came in with.""")

# ─────────────────────────────────────────────────────────────────────────────
# S7 — Level 3: Context Engineer
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, GREEN)
tag(s, "Level 3")
level_badge(s, 3, "CONTEXT ENGINEER", GREEN, Inches(0.5), Inches(0.18))

tb(s, "Level 3: Context Engineer",
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75), size=36, bold=True)

# CLAUDE.md card
c = card(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.2), border=GREEN)
tf = c.text_frame; tf.word_wrap = True
p1 = tf.paragraphs[0]
r1 = p1.add_run(); r1.text = "  CLAUDE.md"
r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = GREEN
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r2 = p2.add_run()
r2.text = ("  Your project's identity card.\n"
           "  Technologies, standards, rules,\n"
           "  and conventions — written once.")
r2.font.size = Pt(14); r2.font.color.rgb = GRAY

blist(s, [
    "Loaded automatically at every session start",
    "Defines: stack, code style, naming, deployment",
    "Claude follows your rules — every single time",
], Inches(0.5), Inches(4.2), Inches(5.8), Inches(1.8), size=16)

# Less is more
c2 = card(s, Inches(6.7), Inches(1.8), Inches(6.2), Inches(2.5), border=RED, fill=BG_CARD2)
tb(s, "CRITICAL RULE", Inches(6.9), Inches(1.95), Inches(5.8), Inches(0.45),
   size=12, bold=True, color=RED)
tb(s, "LESS IS MORE", Inches(6.9), Inches(2.4), Inches(5.8), Inches(0.85),
   size=32, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s, "Bloated context files make models perform WORSE",
   Inches(6.9), Inches(3.2), Inches(5.8), Inches(0.55),
   size=15, color=GRAY, align=PP_ALIGN.CENTER)

# Token decay stat
stat_box(s, "100K", "tokens — decay begins here", Inches(6.8), Inches(4.5), YELLOW)
tb(s, "Even with 1M token window, performance drops after 100K.\nBest practice: summarize → /clear → paste → continue.",
   Inches(6.8), Inches(5.85), Inches(6.2), Inches(0.75),
   size=13, color=GRAY, italic=True)

snum(s, 7)
add_notes(s, """Level 3 is where you start engineering your context — and this is where most developers plateau.
CLAUDE.md is a markdown file in your project root. Claude reads it automatically at the start of every session. It's your project's identity card.
Put in: your tech stack, coding standards, naming conventions, deployment process. Write it once. Claude follows it forever.
Now here is the most counter-intuitive thing I'm going to say today: LESS IS MORE.
Research shows that bloated context files make model performance WORSE — not better. More input tokens cause what's called context decay.
Even with a 1 million token context window, performance starts dropping around 100K tokens.
The power move: keep your CLAUDE.md tight and focused. Summarize the session periodically, use /clear, and paste only the essential context back in. That's how professionals work.""")

# ─────────────────────────────────────────────────────────────────────────────
# S8 — Level 4: Tool Master (MCP)
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, YELLOW)
tag(s, "Level 4")
level_badge(s, 4, "TOOL MASTER", YELLOW, Inches(0.5), Inches(0.18))

tb(s, "Level 4: Tool Master — MCP",
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75), size=36, bold=True)
tb(s, "MCP = Model Context Protocol — connects Claude to external tools.",
   Inches(0.6), Inches(1.65), Inches(12), Inches(0.5), size=19, color=YELLOW)

tools = [
    ("Supabase",    "Database", CYAN),
    ("Figma",       "Design",   GREEN),
    ("Playwright",  "Testing",  ORANGE),
    ("GitHub",      "Git ops",  GRAY),
    ("Slack",       "Comms",    TEAL),
    ("Custom APIs", "Anything", RED),
]
x, y = Inches(0.5), Inches(2.45)
for i, (name, cat, col) in enumerate(tools):
    c = card(s, x, y, Inches(3.8), Inches(0.9), border=col)
    tf = c.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"  {name}  "
    r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = col
    r2 = p.add_run(); r2.text = f"— {cat}"
    r2.font.size = Pt(13); r2.font.color.rgb = GRAY
    if i % 3 == 2:
        x = Inches(0.5); y += Inches(1.05)
    else:
        x += Inches(4.3)

# Candy shop trap
c2 = card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.05), border=RED, fill=BG_CARD2)
tf2 = c2.text_frame
p1 = tf2.paragraphs[0]
r1 = p1.add_run(); r1.text = "  The Candy Shop Trap:  "
r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = RED
r2 = p1.add_run()
r2.text = "Connecting 15 tools = Claude picks the wrong one. More options = worse decisions."
r2.font.size = Pt(14); r2.font.color.rgb = GRAY
p2 = tf2.add_paragraph(); p2.space_before = Pt(4)
r3 = p2.add_run()
r3.text = "  Rule: surgical precision — only connect the tools you actually need for this task."
r3.font.size = Pt(13); r3.font.color.rgb = YELLOW; r3.font.italic = True

snum(s, 8)
add_notes(s, """Level 4 is where Claude Code becomes genuinely powerful — you connect it to external tools via MCP, the Model Context Protocol.
Supabase for your database. Figma for design specs. Playwright for automated testing. GitHub for git operations. The list goes on.
When Claude has these tools, it doesn't just write code — it writes code, runs the database migration, checks the design spec, and runs the tests. All in one session.
BUT — and this is critical — there is a candy shop trap.
If you connect 15 tools, Claude gets confused about which one to use. More options lead to worse decisions. This is a research-backed finding about language model behavior.
Surgical precision is the rule. Only connect what you actually need for the task at hand. Three well-chosen tools beat fifteen vague ones every time.
Also: make sure you understand the basic concepts — frontend, backend, database, auth, API, deploy. Claude can't fill knowledge gaps you haven't defined.""")

# ─────────────────────────────────────────────────────────────────────────────
# S9 — Level 5: Skilled Craftsman
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, ORANGE)
tag(s, "Level 5")
level_badge(s, 5, "SKILLED CRAFTSMAN", ORANGE, Inches(0.5), Inches(0.18))

tb(s, "Level 5: Skilled Craftsman",
   Inches(0.6), Inches(0.85), Inches(12), Inches(0.75), size=36, bold=True)
tb(s, "Create custom skills — teach Claude HOW to do things your way.",
   Inches(0.6), Inches(1.65), Inches(12), Inches(0.55), size=19, color=ORANGE)

blist(s, [
    "Skills = custom MD files that define HOW Claude approaches a task",
    "Use the Skill Creator plugin to build your own",
    "Example: /deploy, /review-pr, /refactor, /write-tests",
    "Your Claude setup becomes unique — nobody else's works like yours",
], Inches(0.5), Inches(2.5), Inches(7.0), Inches(2.8), size=17)

# Quality over quantity
c = card(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.85), border=ORANGE, fill=BG_CARD2)
tf = c.text_frame
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = "  Quality over quantity:  "
r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = ORANGE
r2 = p.add_run(); r2.text = "20 carefully crafted skills > 100 random ones."
r2.font.size = Pt(15); r2.font.color.rgb = GRAY

# Skill cards
skill_examples = [
    ("/deploy",       ORANGE),
    ("/review-pr",    CYAN),
    ("/refactor",     GREEN),
    ("/write-tests",  TEAL),
    ("/analyse-data", YELLOW),
]
x = Inches(7.5)
y_s = Inches(2.5)
for i, (skill, col) in enumerate(skill_examples):
    c = card(s, x, y_s, Inches(2.8), Inches(0.65), border=col)
    tb_s = c.text_frame.paragraphs[0]
    tb_s.alignment = PP_ALIGN.CENTER
    r = tb_s.add_run(); r.text = skill
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = col
    y_s += Inches(0.8)

snum(s, 9)
add_notes(s, """Level 5 is where you start customizing Claude Code to work exactly the way you work.
Skills are markdown files that teach Claude HOW to approach specific tasks — not just what to do, but your workflow, your standards, your style.
For example: a /deploy skill that knows your exact deployment process. A /review-pr skill that checks for the specific things your team cares about. A /write-tests skill that matches your testing framework.
Use the Skill Creator plugin to build these systematically.
Here's the key insight: 20 carefully crafted skills are worth more than 100 generic ones.
Each skill you create makes your Claude setup more unique, more powerful, and more aligned to your actual work.
Nobody else's Claude works like yours. That's the point.""")

# ─────────────────────────────────────────────────────────────────────────────
# S10 — Level 6: Orchestrator
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, RED)
tag(s, "Level 6")
level_badge(s, 6, "ORCHESTRATOR", RED, Inches(0.5), Inches(0.18))

tb(s, "Level 6: Orchestrator — Agent Teams",
   Inches(0.6), Inches(0.85), Inches(12), Inches(0.75), size=34, bold=True)
tb(s, "Run multiple Claude agents in parallel. Assign. Coordinate. Ship.",
   Inches(0.6), Inches(1.62), Inches(12), Inches(0.55), size=18, color=RED)

methods = [
    ("1", "Multi-Terminal",
     "Open multiple terminals, each running its own Claude session. Easy, no coordination needed.",
     GREEN, Inches(0.5)),
    ("2", "Git Worktree / Sub-agents",
     "Claude creates its own sub-agents and assigns them tasks. Medium complexity.",
     YELLOW, Inches(4.5)),
    ("3", "Agent Team",
     "Agents communicate with each other. Supervisor distributes tasks, agents report back. Advanced.",
     RED, Inches(8.5)),
]
for num, title, desc, col, x in methods:
    c = card(s, x, Inches(2.5), Inches(4.3), Inches(2.5), border=col)
    tf = c.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r0 = p1.add_run(); r0.text = f"  Method {num}  "
    r0.font.size = Pt(11); r0.font.color.rgb = col; r0.font.bold = True
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    r1 = p2.add_run(); r1.text = f"  {title}"
    r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = WHITE
    p3 = tf.add_paragraph(); p3.space_before = Pt(6)
    r2 = p3.add_run(); r2.text = f"  {desc}"
    r2.font.size = Pt(12); r2.font.color.rgb = GRAY

# Warning
c2 = card(s, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.9), border=RED, fill=BG_CARD2)
tf2 = c2.text_frame
p = tf2.paragraphs[0]
r1 = p.add_run(); r1.text = "  Warning:  "
r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = RED
r2 = p.add_run()
r2.text = "Agent Teams are high token cost and still experimental. Start with Method 1. Graduate to 2 and 3 as needed."
r2.font.size = Pt(14); r2.font.color.rgb = GRAY

snum(s, 10)
add_notes(s, """Level 6 is where you stop being a single developer and start running a team — except the team is made of Claude agents.
Three methods, in order of complexity.
Method 1 — Multi-terminal: open three terminal windows, each with its own Claude session. One builds the frontend, one builds the backend, one writes the tests. No coordination overhead. Start here.
Method 2 — Git Worktree and sub-agents: Claude creates its own sub-agents and assigns them to isolated branches. It coordinates the work. Medium complexity.
Method 3 — Agent Team: a supervisor agent distributes tasks, sub-agents execute, and agents communicate with each other. This is the most powerful setup and also the most experimental.
Important warning: Agent Teams are expensive on tokens and the tooling is still maturing. Don't start here. Get comfortable with Methods 1 and 2 first.
When you're running Level 6 effectively, your throughput as a developer multiplies by a factor of three or more.""")

# ─────────────────────────────────────────────────────────────────────────────
# S11 — Live Demo: Retail Analytics Project
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, CYAN)
tag(s, "Live Demo")

tb(s, "Live Demo: Retail Analytics Project",
   Inches(0.6), Inches(0.85), Inches(12), Inches(0.75), size=34, bold=True)
tb(s, "Everything you just learned — applied in one real session.",
   Inches(0.6), Inches(1.65), Inches(12), Inches(0.5), size=19, color=CYAN)

# Stats row
stats = [
    ("541,909", "Rows loaded"),
    ("£9.7M",   "Net revenue"),
    ("10,624",  "Cancellations\nflagged"),
    ("6",       "Dashboard\ntabs"),
]
x = Inches(0.5)
for val, lbl in stats:
    stat_box(s, val, lbl, x, Inches(2.3), CYAN)
    x += Inches(3.1)

# What Claude Code did
tb(s, "What Claude Code built — from plain English prompts only:",
   Inches(0.5), Inches(4.1), Inches(12), Inches(0.45),
   size=14, bold=True, color=WHITE)

files = [
    ("explore.py",           "Initial data exploration — columns, dtypes, nulls, date range"),
    ("data_prep.py",         "Cleaning, cancellation flagging, LineTotal / GrossLineTotal / ReturnLineTotal"),
    ("retail_dashboard.py",  "Full Streamlit app — KPIs, charts, map, QA tab — all interactive"),
]
y = Inches(4.65)
for fname, desc in files:
    c = card(s, Inches(0.5), y, Inches(12.3), Inches(0.62), border=CYAN)
    tf = c.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"  {fname}  "
    r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = CYAN
    r2 = p.add_run(); r2.text = f"— {desc}"
    r2.font.size = Pt(13); r2.font.color.rgb = GRAY
    y += Inches(0.72)

snum(s, 11)
add_notes(s, """Now let me show you everything we just talked about in practice — with a real project I built live using Claude Code.
The dataset: 541,909 rows of UK retail transactions. Real data. No clean sample.
I gave Claude Code plain English instructions. No pre-written scripts. No templates.
What came out in one session: three production-ready Python files.
explore.py — Claude explored the data first: columns, data types, null counts, date range. It verified numbers before assuming anything.
data_prep.py — Claude added calculated columns, flagged 10,624 cancellations WITHOUT deleting them, and printed a full data quality summary. It understood business logic from natural language.
retail_dashboard.py — A full Streamlit application: 6 tabs, interactive Plotly charts, a choropleth map, sidebar controls, and a QA tab. Built without a single line of boilerplate written by me.
This is what Levels 2 through 4 look like in practice: planning, context via CLAUDE.md, and surgical tool use.
[Launch the dashboard live if demo environment is ready]""")

# ─────────────────────────────────────────────────────────────────────────────
# S12 — Key Takeaways
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, CYAN)
tag(s, "Takeaways")

tb(s, "Key Takeaways",
   Inches(0.6), Inches(0.85), Inches(12), Inches(0.75), size=38, bold=True)

# Level progression
tb(s, "Each level builds on the previous — don't skip.",
   Inches(0.6), Inches(1.65), Inches(12), Inches(0.5), size=19, color=CYAN)

colors6 = [CYAN, TEAL, GREEN, YELLOW, ORANGE, RED]
labels  = ["L1\nCMD", "L2\nPLAN", "L3\nCTX", "L4\nMCP", "L5\nSKILL", "L6\nAGENT"]
x = Inches(0.6)
for col, lbl in zip(colors6, labels):
    c = s.shapes.add_textbox(x, Inches(2.3), Inches(1.7), Inches(0.75))
    c.fill.solid(); c.fill.fore_color.rgb = col
    c.line.fill.background()
    tf = c.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lbl
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BG
    x += Inches(1.95)

# 3 golden rules
rules = [
    ("1", "Never give commands without context", "No context = average output. Always describe what you want and why.", CYAN),
    ("2", "Less is more", "Context files, tools, skills — keep them tight and purposeful.", GREEN),
    ("3", "Treat Claude as a partner, not a tool", "Collaborate. Plan. Challenge ideas. You'll get far better results.", ORANGE),
]
y = Inches(3.3)
for num, rule, detail, col in rules:
    c = card(s, Inches(0.5), y, Inches(12.3), Inches(0.9), border=col)
    tf = c.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"  {num}.  {rule}  — "
    r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = col
    r2 = p.add_run(); r2.text = detail
    r2.font.size = Pt(14); r2.font.color.rgb = GRAY
    y += Inches(1.02)

tb(s, '"The winners of tomorrow are learning Claude Code today."',
   Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.65),
   size=17, color=YELLOW, italic=True, align=PP_ALIGN.CENTER)

snum(s, 12)
add_notes(s, """Let me bring everything together.
Six levels. Each one builds on the previous. You cannot skip levels and get consistent results.
Level 1: give context. Level 2: plan before you build. Level 3: keep context tight. Level 4: use tools surgically. Level 5: build your own skills. Level 6: run agent teams.
Three golden rules to take home:
One — never give commands without context. No context equals average output. This is the most common mistake and the easiest to fix.
Two — less is more. In context, in tools, in skills. Restraint improves performance.
Three — treat Claude as a partner. Plan with it. Challenge it. Ask it what it would do differently. The collaborative mindset produces far better outcomes than the command mindset.
And I'll leave you with this: the winners of tomorrow are learning Claude Code today.
Not because it's a trend. Because it fundamentally changes what one person can build.""")

# ─────────────────────────────────────────────────────────────────────────────
# S13 — Thank You / Q&A
# ─────────────────────────────────────────────────────────────────────────────
s = slide()
bar(s, CYAN)

tb(s, "Thank You", Inches(0.7), Inches(1.2), W - Inches(1.4), Inches(1.4),
   size=64, bold=True, align=PP_ALIGN.CENTER)

tb(s, "Questions?", Inches(0.7), Inches(2.8), W - Inches(1.4), Inches(0.8),
   size=36, color=CYAN, align=PP_ALIGN.CENTER)

next_steps = [
    "anthropic.com  —  Get your API key",
    "npm install -g @anthropic-ai/claude-code  —  Install",
    "youtube.com/watch?v=s2KVJirRoRs  —  Watch the deep dive",
    "Start with your own messy dataset or codebase",
]
blist(s, next_steps, Inches(3.0), Inches(4.0), Inches(7.5), Inches(2.2),
      size=16, color=GRAY, marker="  ->  ", gap=8)

tb(s, '"The only difference between you and someone already automating with Claude\nis one terminal command — and now you know which one to type."',
   Inches(0.7), Inches(6.2), W - Inches(1.4), Inches(0.9),
   size=15, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

snum(s, 13)
add_notes(s, """Thank you all for your time and attention.
Let me leave you with the only thing you actually need to do after this talk:
Go to anthropic.com. Get an API key. Install Claude Code. Open your project folder. Type 'claude'. Describe what you want.
That's it. That is the entire on-ramp.
Watch the YouTube link if you want the deep technical architecture walkthrough — it's the best single resource I've found.
If you want to see the retail analytics dashboard code we demoed today, the repo is available — ask me after.
The floor is open for questions. Thank you.""")

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = "output/Claude_Code_Tech_Talk_v4.pptx"
prs.save(out)
print(f"Saved: {out}  ({TOTAL} slides, speaker notes on every slide)")
